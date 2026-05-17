import os
import sys

# Ensure python-pptx is installed or try installing it
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not installed. Installing it now...")
    os.system(sys.executable + " -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Color Palettes
    BG_COLOR = RGBColor(15, 18, 29)       # Deep slate (#0f121d)
    CARD_BG = RGBColor(22, 28, 45)        # Dark navy card (#161c2d)
    BORDER_COLOR = RGBColor(51, 65, 85)   # Border gray (#334155)
    
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_LIGHT = RGBColor(209, 213, 219)
    TEXT_MUTED = RGBColor(156, 163, 175)
    
    ACCENT_INDIGO = RGBColor(129, 140, 248) # Indigo glow
    ACCENT_CYAN = RGBColor(34, 211, 238)     # Cyan focus
    ACCENT_EMERALD = RGBColor(52, 211, 153)  # Emerald focus
    ACCENT_AMBER = RGBColor(251, 191, 36)    # Amber accent

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def apply_dark_background(slide):
        # Draw background rectangle
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG_COLOR
        rect.line.fill.background() # No border
        return slide

    def add_title(slide, text, color=TEXT_WHITE):
        tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Outfit"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

    def add_footer(slide, current, total=12):
        # Footer branding
        brand_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.9), Inches(5.0), Inches(0.3))
        tf_brand = brand_box.text_frame
        p_brand = tf_brand.paragraphs[0]
        p_brand.text = "TomoClub | Curriculum Design Proposal"
        p_brand.font.name = "Inter"
        p_brand.font.size = Pt(10)
        p_brand.font.color.rgb = TEXT_MUTED
        
        # Slide counter
        counter_box = slide.shapes.add_textbox(Inches(11.583), Inches(6.9), Inches(1.0), Inches(0.3))
        tf_counter = counter_box.text_frame
        p_counter = tf_counter.paragraphs[0]
        p_counter.alignment = PP_ALIGN.RIGHT
        p_counter.text = f"{current} / {total}"
        p_counter.font.name = "Outfit"
        p_counter.font.size = Pt(12)
        p_counter.font.bold = True
        p_counter.font.color.rgb = ACCENT_INDIGO

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide1)
    
    # Decorative colored glow block
    glow = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(6.0), Inches(0.08))
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_INDIGO
    glow.line.fill.background()

    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Designing Smarter Classrooms:"
    p1.font.name = "Outfit"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "AI Use in Assessments & Personalisation"
    p2.font.name = "Outfit"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN
    
    # Metadata Box
    meta_box = slide1.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(8.0), Inches(1.5))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    p3 = tf_meta.paragraphs[0]
    p3.text = "Session Flow & Curriculum Design Proposal for K-12 Educators"
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_LIGHT
    p3.space_after = Pt(20)
    
    p4 = tf_meta.add_paragraph()
    p4.text = "Prepared for: TomoClub Selection Committee    |    By: Pramit Das"
    p4.font.name = "Inter"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_INDIGO
    
    add_footer(slide1, 1)

    # ==========================================
    # SLIDE 2: CONTEXT & PAIN POINTS (TRILEMMA)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide2)
    add_title(slide2, "The K-12 Educator Trilemma")
    
    # Subtitle
    sub_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.4))
    sub_box.text_frame.paragraphs[0].text = "Why educators struggle with assessment and differentiation in the AI era"
    sub_box.text_frame.paragraphs[0].font.name = "Inter"
    sub_box.text_frame.paragraphs[0].font.size = Pt(14)
    sub_box.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # Three Grid Cards
    cards_data = [
        ("🚨 Severe Workload Burnout", "Teachers spend upwards of 10+ hours a week grading, designing rubrics, and lesson planning, leaving them exhausted and reducing active face-to-face classroom relationship building.", ACCENT_INDIGO),
        ("🧩 The Personalisation Gap", "With over 30+ students per K-12 classroom, manual differentiation is virtually impossible. Accommodating diverse profiles remains a paper checklist rather than active classroom reality.", ACCENT_CYAN),
        ("🛡️ The Academic Integrity Crisis", "Plagiarism software produces false positives and is easily bypassed. Trying to lock students out of AI locks schools into a surveillance battle, rather than process grading.", ACCENT_AMBER)
    ]
    
    card_width = Inches(3.68)
    card_height = Inches(4.3)
    start_left = Inches(0.75)
    gap = Inches(0.4)
    top_pos = Inches(2.0)
    
    for idx, (title, desc, color) in enumerate(cards_data):
        current_left = start_left + idx * (card_width + gap)
        
        # Background shape
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Color Accent border top line
        acc = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_left, top_pos, card_width, Inches(0.08))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()
        
        # Content box
        t_box = slide2.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), card_width - Inches(0.5), card_height - Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Outfit"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(14)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Inter"
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.line_spacing = 1.3
        
    add_footer(slide2, 2)

    # ==========================================
    # SLIDE 3: WORKSHOP OVERVIEW & PARAMETERS
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide3)
    add_title(slide3, "Workshop Overview & Parameters")
    
    # Left Column: Table Parameters
    param_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    param_card.fill.solid()
    param_card.fill.fore_color.rgb = CARD_BG
    param_card.line.color.rgb = BORDER_COLOR
    
    # Title
    t_box_l = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_l = t_box_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "Session Core Specs"
    p_l.font.name = "Outfit"
    p_l.font.size = Pt(20)
    p_l.font.bold = True
    p_l.font.color.rgb = ACCENT_CYAN
    p_l.space_after = Pt(15)
    
    specs = [
        ("Target Audience", "K-12 Educators, Instructional Leads, Administrators"),
        ("Session Duration", "120 Minutes (2 Hours) - Perfect for active sandboxes"),
        ("Delivery Style", "Blended / Collaborative (Max Hands-on Activity)"),
        ("Key Pedagogy", "The Sandwich Model: Action -> Theory -> Action -> Review"),
        ("Key Resource", "Classroom Prompt Sandbox and Assessment Folder")
    ]
    
    for label, val in specs:
        p_label = tf_l.add_paragraph()
        p_label.text = f"• {label}:"
        p_label.font.name = "Inter"
        p_label.font.size = Pt(13)
        p_label.font.bold = True
        p_label.font.color.rgb = TEXT_WHITE
        p_label.space_after = Pt(2)
        
        p_val = tf_l.add_paragraph()
        p_val.text = f"   {val}"
        p_val.font.name = "Inter"
        p_val.font.size = Pt(12)
        p_val.font.color.rgb = TEXT_LIGHT
        p_val.space_after = Pt(8)

    # Right Column: Strategy Description
    strat_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    strat_card.fill.solid()
    strat_card.fill.fore_color.rgb = CARD_BG
    strat_card.line.color.rgb = BORDER_COLOR
    
    t_box_r = slide3.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_r = t_box_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "Why This Session Model Works"
    p_r.font.name = "Outfit"
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_INDIGO
    p_r.space_after = Pt(15)
    
    strats = [
        ("Active Engagement over Passive Listening", "Traditional PD lectures fail to translate into classroom practice. We keep direct slide presentation under 20 minutes, allocating the remaining 100 minutes to active creation, prompt testing, and peer auditing."),
        ("Direct Utility & Classroom-Ready Outcomes", "Every teacher leaves the session with at least one custom, standards-aligned assessment prompt and a multi-level differentiated reading package ready for their students next morning."),
        ("Combines Technological & Critical Literacy", "Rather than just teaching 'how to copy prompts', we emphasize 'Critical AI Auditing', showing teachers how to check outputs for hallucinations and level adjustments.")
    ]
    
    for title, desc in strats:
        p_st = tf_r.add_paragraph()
        p_st.text = f"⚡ {title}"
        p_st.font.name = "Inter"
        p_st.font.size = Pt(13)
        p_st.font.bold = True
        p_st.font.color.rgb = TEXT_WHITE
        p_st.space_after = Pt(2)
        
        p_sd = tf_r.add_paragraph()
        p_sd.text = desc
        p_sd.font.name = "Inter"
        p_sd.font.size = Pt(11)
        p_sd.font.color.rgb = TEXT_LIGHT
        p_sd.space_after = Pt(10)
        p_sd.line_spacing = 1.2

    add_footer(slide3, 3)

    # ==========================================
    # SLIDE 4: LEARNING OBJECTIVES
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide4)
    add_title(slide4, "Learning Objectives (Bloom's Taxonomy)")
    
    cards_obj = [
        ("1. UNDERSTAND", "Shift in Assessment Philosophy", "Educators will analyze the structural shift from grading 'product' (essays, static tests) to grading 'process' (concept outlines, oral defenses, local prompt adaptations), allowing them to design robust, AI-resistant assessment strategies.", ACCENT_INDIGO),
        ("2. APPLY", "Direct Prompt Sandbox Design", "Educators will construct high-fidelity, role-based AI prompt templates that automatically generate formative quizzes, multi-level reading passages, and customized grading rubrics aligned directly with grade-appropriate standards.", ACCENT_CYAN),
        ("3. EVALUATE", "Critical AI Output Auditing", "Educators will critique AI-generated outputs for potential hallucinations, grade-level vocabulary mismatches, and hidden societal biases, mastering the core skill of human-in-the-loop validation.", ACCENT_EMERALD)
    ]
    
    for idx, (head, sub, desc, color) in enumerate(cards_obj):
        current_left = start_left + idx * (card_width + gap)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        acc = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_left, top_pos, card_width, Inches(0.08))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()
        
        t_box = slide4.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), card_width - Inches(0.5), card_height - Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = "Outfit"
        p_sub.font.size = Pt(16)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Inter"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT
        p2.line_spacing = 1.3
        
    add_footer(slide4, 4)

    # ==========================================
    # SLIDE 5: WORKSHOP FLOW (120-MINS TIMELINE)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide5)
    add_title(slide5, "Workshop Session Flow (120 Mins)")
    
    # Draw horizontal timeline track
    track = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.6), Inches(11.83), Inches(0.04))
    track.fill.solid()
    track.fill.fore_color.rgb = BORDER_COLOR
    track.line.fill.background()
    
    # 6 Steps
    steps = [
        ("01", "15 Mins", "The Hook & Perceptions", "Pre-session Mentimeter checking comfort levels.", ACCENT_INDIGO, Inches(0.75)),
        ("02", "20 Mins", "Mini-Input Prompting", "Pedagogical role, context, and standard alignment.", ACCENT_CYAN, Inches(2.81)),
        ("03", "30 Mins", "Assessment Sandbox", "Cooperative grade breakout prompt engineering.", ACCENT_EMERALD, Inches(4.87)),
        ("04", "25 Mins", "Personalisation Sandbox", "Scaffolding same science topic for 3 student pathways.", ACCENT_AMBER, Inches(6.93)),
        ("05", "15 Mins", "Integrity & Safety Checklist", "Process grading over surveillance and data privacy.", ACCENT_INDIGO, Inches(8.99)),
        ("06", "15 Mins", "Action & Reflections", "Google Exit Ticket, async portfolio review.", ACCENT_CYAN, Inches(11.05))
    ]
    
    for num, duration, title, desc, color, left_pos in steps:
        # Step node (circle)
        node = slide5.shapes.add_shape(MSO_SHAPE.OVAL, left_pos, Inches(3.42), Inches(0.4), Inches(0.4))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        
        # Step Number Text inside circle
        p_node = node.text_frame.paragraphs[0]
        p_node.text = num
        p_node.font.name = "Outfit"
        p_node.font.size = Pt(11)
        p_node.font.bold = True
        p_node.font.color.rgb = BG_COLOR
        p_node.alignment = PP_ALIGN.CENTER
        
        # Text block above timeline (Duration & Title)
        above_box = slide5.shapes.add_textbox(left_pos - Inches(0.6), Inches(1.8), Inches(1.6), Inches(1.5))
        tf_a = above_box.text_frame
        tf_a.word_wrap = True
        
        pa_dur = tf_a.paragraphs[0]
        pa_dur.text = duration
        pa_dur.font.name = "Outfit"
        pa_dur.font.size = Pt(14)
        pa_dur.font.bold = True
        pa_dur.font.color.rgb = color
        pa_dur.alignment = PP_ALIGN.CENTER
        
        pa_tit = tf_a.add_paragraph()
        pa_tit.text = title
        pa_tit.font.name = "Outfit"
        pa_tit.font.size = Pt(12)
        pa_tit.font.bold = True
        pa_tit.font.color.rgb = TEXT_WHITE
        pa_tit.alignment = PP_ALIGN.CENTER
        
        # Text block below timeline (Description)
        below_box = slide5.shapes.add_textbox(left_pos - Inches(0.8), Inches(4.0), Inches(2.0), Inches(2.2))
        tf_b = below_box.text_frame
        tf_b.word_wrap = True
        
        pb_desc = tf_b.paragraphs[0]
        pb_desc.text = desc
        pb_desc.font.name = "Inter"
        pb_desc.font.size = Pt(10)
        pb_desc.font.color.rgb = TEXT_LIGHT
        pb_desc.alignment = PP_ALIGN.CENTER
        pb_desc.line_spacing = 1.2
        
    add_footer(slide5, 5)

    # ==========================================
    # SLIDE 6: ACTIVITY 1 - ASSESSMENT SANDBOX
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide6)
    add_title(slide6, "Sandbox 1: Assessment Prompt Battle")
    
    # Left Box: Activity details
    act_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    act_card.fill.solid()
    act_card.fill.fore_color.rgb = CARD_BG
    act_card.line.color.rgb = BORDER_COLOR
    
    t_box_act = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_act = t_box_act.text_frame
    tf_act.word_wrap = True
    
    p = tf_act.paragraphs[0]
    p.text = "Cooperative Breakout Mechanics"
    p.font.name = "Outfit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(12)
    
    steps_act = [
        ("Objective", "Design a custom quiz or assessment rubric targeting standard-aligned topics."),
        ("Interactive Sandbox", "Teachers use the dashboard builder, changing Grade, Subject, and Format dropdowns to dynamically compile role-based prompts."),
        ("The 'Battle' Challenge", "Groups copy the compiled prompt, generate the quiz in Claude/ChatGPT, and must perform a 'Critical AI Audit'—making at least 3 manual edits to the output to fix grade levels, bias, or wording."),
        ("Collaborative Sharing", "Polished rubrics are uploaded to Padlet for live critique by the facilitator and other groups.")
    ]
    
    for head, desc in steps_act:
        p_head = tf_act.add_paragraph()
        p_head.text = f"🎯 {head}:"
        p_head.font.name = "Inter"
        p_head.font.size = Pt(12)
        p_head.font.bold = True
        p_head.font.color.rgb = TEXT_WHITE
        p_head.space_after = Pt(1)
        
        p_desc = tf_act.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Inter"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_LIGHT
        p_desc.space_after = Pt(8)

    # Right Box: Dynamic Prompt Example
    pr_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    pr_card.fill.solid()
    pr_card.fill.fore_color.rgb = CARD_BG
    pr_card.line.color.rgb = BORDER_COLOR
    
    t_box_pr = slide6.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_pr = t_box_pr.text_frame
    tf_pr.word_wrap = True
    
    p_pr = tf_pr.paragraphs[0]
    p_pr.text = "Role-Based AI Prompt Template"
    p_pr.font.name = "Outfit"
    p_pr.font.size = Pt(20)
    p_pr.font.bold = True
    p_pr.font.color.rgb = ACCENT_CYAN
    p_pr.space_after = Pt(15)
    
    prompt_text = (
        "\"Act as an expert K-12 curriculum designer. Create a standards-aligned Grade 8 assessment for Science: Photosynthesis.\n\n"
        "The output must be a 5-question multiple choice quiz with detailed rubrics. Ensure the reading level is age-appropriate, vocabulary is aligned with Next Generation Science Standards (NGSS), and include common misconceptions in distractor choices.\n\n"
        "Do not include generic questions; make all distractors highly plausible.\""
    )
    
    p_body = tf_pr.add_paragraph()
    p_body.text = prompt_text
    p_body.font.name = "Inter"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = ACCENT_EMERALD
    p_body.font.italic = True
    p_body.line_spacing = 1.3
    
    p_ref = tf_pr.add_paragraph()
    p_ref.text = "\n💡 Refinement Checklist:\n1. Did the AI define roles and standard limits?\n2. Did you test for distractor plausibility?"
    p_ref.font.name = "Inter"
    p_ref.font.size = Pt(12)
    p_ref.font.color.rgb = TEXT_LIGHT
    
    add_footer(slide6, 6)

    # ==========================================
    # SLIDE 7: ACTIVITY 2 - PERSONALISATION SANDBOX
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide7)
    add_title(slide7, "Sandbox 2: Differentiating a Core Concept")
    
    # Left side: Profiles
    prof_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    prof_card.fill.solid()
    prof_card.fill.fore_color.rgb = CARD_BG
    prof_card.line.color.rgb = BORDER_COLOR
    
    t_box_prof = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_prof = t_box_prof.text_frame
    tf_prof.word_wrap = True
    
    p_prof = tf_prof.paragraphs[0]
    p_prof.text = "The Differentiated Classroom Challenge"
    p_prof.font.name = "Outfit"
    p_prof.font.size = Pt(20)
    p_prof.font.bold = True
    p_prof.font.color.rgb = ACCENT_CYAN
    p_prof.space_after = Pt(12)
    
    p_goal = tf_prof.add_paragraph()
    p_goal.text = "Goal: Differentiate a single core concept (e.g. Ecosystem Disturbances) for three distinct student profiles:"
    p_goal.font.name = "Inter"
    p_goal.font.size = Pt(12)
    p_goal.font.bold = True
    p_goal.font.color.rgb = TEXT_WHITE
    p_goal.space_after = Pt(10)
    
    profs = [
        ("Profile A: Struggling Reader (Lexile 600L)", "Needs simplified vocabulary, shorter sentences, and direct visual/analogical scaffolding to grasp core concepts.", ACCENT_INDIGO),
        ("Profile B: Multilingual Learner (ELL Level 2)", "Needs dual-language vocabulary list, sentence starters, and bulleted process steps to show scientific understanding.", ACCENT_EMERALD),
        ("Profile C: Advanced / Gifted Student", "Needs high-rigor inquiry prompt exploring open mathematical modeling of ecosystem disruptions and self-guided extension.", ACCENT_AMBER)
    ]
    
    for title, desc, color in profs:
        p_pt = tf_prof.add_paragraph()
        p_pt.text = f"• {title}"
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(12)
        p_pt.font.bold = True
        p_pt.font.color.rgb = color
        p_pt.space_after = Pt(1)
        
        p_pd = tf_prof.add_paragraph()
        p_pd.text = f"   {desc}"
        p_pd.font.name = "Inter"
        p_pd.font.size = Pt(11)
        p_pd.font.color.rgb = TEXT_LIGHT
        p_pd.space_after = Pt(6)

    # Right side: AI Prompt Template
    diff_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    diff_card.fill.solid()
    diff_card.fill.fore_color.rgb = CARD_BG
    diff_card.line.color.rgb = BORDER_COLOR
    
    t_box_diff = slide7.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_diff = t_box_diff.text_frame
    tf_diff.word_wrap = True
    
    p_diff = tf_diff.paragraphs[0]
    p_diff.text = "AI Prompt Template for Differentiation"
    p_diff.font.name = "Outfit"
    p_diff.font.size = Pt(20)
    p_diff.font.bold = True
    p_diff.font.color.rgb = ACCENT_INDIGO
    p_diff.space_after = Pt(15)
    
    d_prompt = (
        "\"Differentiate the science topic 'Ecosystem Disturbances' for Grade 7. Generate three versions of the reading passage:\n\n"
        "1. Simplified vocabulary at 600L Lexile with graphic organizer scaffold.\n"
        "2. Language scaffolded with bold academic key terms and sentence starters.\n"
        "3. High-rigor extension prompt analyzing competitive species dynamics.\""
    )
    
    p_db = tf_diff.add_paragraph()
    p_db.text = d_prompt
    p_db.font.name = "Inter"
    p_db.font.size = Pt(11)
    p_db.font.color.rgb = ACCENT_EMERALD
    p_db.font.italic = True
    p_db.line_spacing = 1.3
    p_db.space_after = Pt(15)
    
    p_ref_c = tf_diff.add_paragraph()
    p_ref_c.text = "⚡ Crucial Reflective Question:\n\"How do we use AI to remove barriers without removing the productive cognitive struggle required for real learning?\""
    p_ref_c.font.name = "Inter"
    p_ref_c.font.size = Pt(11)
    p_ref_c.font.bold = True
    p_ref_c.font.color.rgb = ACCENT_AMBER
    p_ref_c.line_spacing = 1.2
    
    add_footer(slide7, 7)

    # ==========================================
    # SLIDE 8: ACADEMIC INTEGRITY & RESP. AI
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide8)
    add_title(slide8, "Academic Integrity & Responsible AI Use")
    
    # Two main blocks
    left_c = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    left_c.fill.solid()
    left_c.fill.fore_color.rgb = CARD_BG
    left_c.line.color.rgb = BORDER_COLOR
    
    t_box_l = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_l = t_box_l.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Pedagogical Workarounds (Process > Product)"
    p.font.name = "Outfit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(15)
    
    workarounds = [
        ("Ditch Plagiarism Detectors", "AI writing detectors are easily bypassed by paraphrasing and exhibit heavy bias against non-native speakers, leading to false cheating accusations."),
        ("Grade the Process Path", "Integrate structured checks: assess brainstorming outlines, handwritten drafts, research logs, and require short 2-minute oral defenses in class."),
        ("Hyper-Local Assessments", "Align questions to unique classroom dialogues, field trips, local community events, or specific analog artifacts that AI engines cannot access.")
    ]
    
    for head, desc in workarounds:
        p_h = tf_l.add_paragraph()
        p_h.text = f"🛡️ {head}"
        p_h.font.name = "Inter"
        p_h.font.size = Pt(12)
        p_h.font.bold = True
        p_h.font.color.rgb = TEXT_WHITE
        p_h.space_after = Pt(1)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Inter"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_LIGHT
        p_d.space_after = Pt(10)
        p_d.line_spacing = 1.2

    right_c = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    right_c.fill.solid()
    right_c.fill.fore_color.rgb = CARD_BG
    right_c.line.color.rgb = BORDER_COLOR
    
    t_box_r = slide8.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_r = t_box_r.text_frame
    tf_r.word_wrap = True
    
    p_r = tf_r.paragraphs[0]
    p_r.text = "Responsible AI Safety Checklist"
    p_r.font.name = "Outfit"
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_AMBER
    p_r.space_after = Pt(15)
    
    safety = [
        ("Absolute Student Data Privacy", "Never upload student PII (names, emails, grades, student work) to open generative AI tools. Keep student records secure and offline."),
        ("Human-in-the-Loop Verification", "AI is a co-pilot, not the pilot. Always audit and review generated rubrics or quizzes for accuracy, bias, and grade level before class."),
        ("Transparent Student Boundaries", "Establish clear guidelines. Let students know exactly when AI use is authorized (e.g. brainstorming) vs. when it is prohibited (e.g. core thinking).")
    ]
    
    for head, desc in safety:
        p_h = tf_r.add_paragraph()
        p_h.text = f"🔒 {head}"
        p_h.font.name = "Inter"
        p_h.font.size = Pt(12)
        p_h.font.bold = True
        p_h.font.color.rgb = TEXT_WHITE
        p_h.space_after = Pt(1)
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Inter"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_LIGHT
        p_d.space_after = Pt(10)
        p_d.line_spacing = 1.2

    add_footer(slide8, 8)

    # ==========================================
    # SLIDE 9: MEL FRAMEWORK - SYNC/LIVE
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide9)
    add_title(slide9, "Monitoring, Evaluation & Learning (Live)")
    
    cards_mel = [
        ("Pre-Session Diagnostic", "Mentimeter Perceptions", "Conducted in first 15 mins. Gauges baseline comfort levels and specific concerns regarding assessment safety. Facilitator adjusts session focus in real-time based on live data.", ACCENT_INDIGO),
        ("In-Session Deliverables", "Facilitator Review Rubric", "During Sandbox breakouts, the facilitator reviews the live collaborative Padlet using a 3-star checklist:\n1. Standard limitation?\n2. Defined role?\n3. Audited corrections?", ACCENT_CYAN),
        ("Post-Session Exit Ticket", "Learning Transfer Metrics", "Conducted via Google Forms in the final 15 minutes. Evaluates self-reported confidence shifts and objective mastery. Target completion rate: 85%+ with average confidence gain of 1.5+ points.", ACCENT_EMERALD)
    ]
    
    for idx, (head, sub, desc, color) in enumerate(cards_mel):
        current_left = start_left + idx * (card_width + gap)
        
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        acc = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_left, top_pos, card_width, Inches(0.08))
        acc.fill.solid()
        acc.fill.fore_color.rgb = color
        acc.line.fill.background()
        
        t_box = slide9.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), card_width - Inches(0.5), card_height - Inches(0.6))
        tf = t_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = "Outfit"
        p_sub.font.size = Pt(16)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Inter"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT
        p2.line_spacing = 1.3
        
    add_footer(slide9, 9)

    # ==========================================
    # SLIDE 10: MEL FRAMEWORK - ASYNC & IMPACT
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide10)
    add_title(slide10, "MEL: Async Portfolios & Long-Term Impact")
    
    # Left side: Portfolio Requirement
    port_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    port_card.fill.solid()
    port_card.fill.fore_color.rgb = CARD_BG
    port_card.line.color.rgb = BORDER_COLOR
    
    t_box_p = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_p = t_box_p.text_frame
    tf_p.word_wrap = True
    
    p = tf_p.paragraphs[0]
    p.text = "Classroom AI Adaptation Folder"
    p.font.name = "Outfit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD
    p.space_after = Pt(12)
    
    p_desc = tf_p.add_paragraph()
    p_desc.text = "To guarantee learning transfer and secure PD certification, educators submit a portfolio folder within 14 days of the session containing:"
    p_desc.font.name = "Inter"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = TEXT_LIGHT
    p_desc.space_after = Pt(10)
    
    items = [
        ("One Co-Created Assessment", "A standards-aligned quiz or rubric, with explicit highlights showing where human-in-the-loop corrections were made."),
        ("One Differentiated Accommodating Set", "A core classroom text or concept scaffolded into three distinct versions corresponding to our student profiles."),
        ("A Critical Reflection Prompt Log", "Reflective log documenting: What worked, what failed, how the AI output was modified, and student feedback on the assessment.")
    ]
    
    for head, text in items:
        p_ih = tf_p.add_paragraph()
        p_ih.text = f"• {head}"
        p_ih.font.name = "Inter"
        p_ih.font.size = Pt(12)
        p_ih.font.bold = True
        p_ih.font.color.rgb = TEXT_WHITE
        p_ih.space_after = Pt(1)
        
        p_it = tf_p.add_paragraph()
        p_it.text = f"   {text}"
        p_it.font.name = "Inter"
        p_it.font.size = Pt(11)
        p_it.font.color.rgb = TEXT_LIGHT
        p_it.space_after = Pt(6)

    # Right side: Long-Term Impact Metrics
    imp_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    imp_card.fill.solid()
    imp_card.fill.fore_color.rgb = CARD_BG
    imp_card.line.color.rgb = BORDER_COLOR
    
    t_box_i = slide10.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_i = t_box_i.text_frame
    tf_i.word_wrap = True
    
    p_i = tf_i.paragraphs[0]
    p_i.text = "90-Day Success & Impact Metrics"
    p_i.font.name = "Outfit"
    p_i.font.size = Pt(20)
    p_i.font.bold = True
    p_i.font.color.rgb = ACCENT_CYAN
    p_i.space_after = Pt(15)
    
    metrics = [
        ("Workload Reduction Index", "Target: Over 3+ hours saved per week on grading and planning. Measured via 30 and 90-day survey diagnostics.", ACCENT_INDIGO),
        ("Personalisation Adoption Rate", "Target: Over 70% of participating teachers actively integrating differentiated pathways monthly in core instruction.", ACCENT_EMERALD),
        ("Process-Oriented Grading Shift", "Qualitative review of school portfolio submissions showing a 40%+ increase in process-oriented grading rubrics.", ACCENT_AMBER)
    ]
    
    for title, desc, color in metrics:
        p_mh = tf_i.add_paragraph()
        p_mh.text = f"📈 {title}"
        p_mh.font.name = "Inter"
        p_mh.font.size = Pt(12)
        p_mh.font.bold = True
        p_mh.font.color.rgb = color
        p_mh.space_after = Pt(1)
        
        p_md = tf_i.add_paragraph()
        p_md.text = desc
        p_md.font.name = "Inter"
        p_md.font.size = Pt(11)
        p_md.font.color.rgb = TEXT_LIGHT
        p_md.space_after = Pt(8)
        p_md.line_spacing = 1.2

    add_footer(slide10, 10)

    # ==========================================
    # SLIDE 11: FORMAT ADAPTATIONS (SYNC VS ASYNC)
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide11)
    add_title(slide11, "Sync vs. Async Format Adaptations")
    
    left_f = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.6), Inches(4.7))
    left_f.fill.solid()
    left_f.fill.fore_color.rgb = CARD_BG
    left_f.line.color.rgb = BORDER_COLOR
    left_f.line.width = Pt(2.0)
    left_f.line.color.rgb = ACCENT_INDIGO
    
    t_box_lf = slide11.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_lf = t_box_lf.text_frame
    tf_lf.word_wrap = True
    
    p = tf_lf.paragraphs[0]
    p.text = "Synchronous Format (Live)"
    p.font.name = "Outfit"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO
    p.space_after = Pt(15)
    
    sync_points = [
        "Live Breakout Rooms: Grade-level cooperative groups work in real-time, screen-sharing prompt builders.",
        "Real-Time Chat & Checks: Facilitator uses interactive chat checks and Mentimeter diagnostic polling.",
        "Collaborative Prompt Battle: Breakouts compete to build and audit the most standards-aligned rubrics on Padlet.",
        "Immediate Peer Critique: Group representatives present prompts, followed by active feedback sessions."
    ]
    
    for pt in sync_points:
        p_pt = tf_lf.add_paragraph()
        p_pt.text = f"👥 {pt}"
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_LIGHT
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.2

    right_f = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.98), Inches(1.6), Inches(5.6), Inches(4.7))
    right_f.fill.solid()
    right_f.fill.fore_color.rgb = CARD_BG
    right_f.line.color.rgb = BORDER_COLOR
    right_f.line.width = Pt(2.0)
    right_f.line.color.rgb = ACCENT_CYAN
    
    t_box_rf = slide11.shapes.add_textbox(Inches(7.23), Inches(1.8), Inches(5.1), Inches(4.2))
    tf_rf = t_box_rf.text_frame
    tf_rf.word_wrap = True
    
    p_rf = tf_rf.paragraphs[0]
    p_rf.text = "Asynchronous Format (Self-Paced)"
    p_rf.font.name = "Outfit"
    p_rf.font.size = Pt(20)
    p_rf.font.bold = True
    p_rf.font.color.rgb = ACCENT_CYAN
    p_rf.space_after = Pt(15)
    
    async_points = [
        "6 Modular Micro-Videos: 120 minutes broken into 6 short interactive videos with embedded check-in checks using Edpuzzle.",
        "Individual Sandbox Builders: Teachers use self-paced prompt worksheets to compile custom prompts locally.",
        "Peer-Review Discussion Boards: Group breakouts are replaced by discussion boards where educators peer-review two colleagues' prompts.",
        "Standards Aligned Portfolio: Both formats submit the exact same adaptation portfolio to achieve PD certification."
    ]
    
    for pt in async_points:
        p_pt = tf_rf.add_paragraph()
        p_pt.text = f"💻 {pt}"
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_LIGHT
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.2

    add_footer(slide11, 11)

    # ==========================================
    # SLIDE 12: CONCLUSION & Q&A
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    apply_dark_background(slide12)
    
    # Decorative glow block
    glow12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(6.0), Inches(0.08))
    glow12.fill.solid()
    glow12.fill.fore_color.rgb = ACCENT_CYAN
    glow12.line.fill.background()

    # Title Box
    title_box12 = slide12.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.83), Inches(2.5))
    tf12 = title_box12.text_frame
    tf12.word_wrap = True
    
    p1 = tf12.paragraphs[0]
    p1.text = "Empowering Teachers."
    p1.font.name = "Outfit"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(10)
    
    p2 = tf12.add_paragraph()
    p2.text = "Designing Future-Ready Classrooms."
    p2.font.name = "Outfit"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_INDIGO
    
    # Call to action text
    cta_box = slide12.shapes.add_textbox(Inches(0.75), Inches(4.6), Inches(10.0), Inches(1.8))
    tf_cta = cta_box.text_frame
    tf_cta.word_wrap = True
    
    p3 = tf_cta.paragraphs[0]
    p3.text = "\"By positioning AI as a collaborative, ethical co-pilot, we reduce teacher burnout, make individual personalization accessible, and allow educators to focus on what they do best: building meaningful human relationships.\""
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_LIGHT
    p3.space_after = Pt(25)
    
    p4 = tf_cta.add_paragraph()
    p4.text = "Join us in reshaping professional development.    |    Q&A Session"
    p4.font.name = "Inter"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_AMBER
    
    add_footer(slide12, 12)

    # Save presentation
    output_filename = "/home/pramit/Desktop/Projects/Pitch2/TomoClub_AI_Assessments_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully to: {output_filename}")

if __name__ == "__main__":
    create_presentation()
